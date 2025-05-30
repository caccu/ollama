from flask import Flask, request, jsonify, send_from_directory, Response, session, current_app
from flask_session import Session
import os
from werkzeug.utils import secure_filename
import tempfile
import sys
import ollama

app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'supersecretkey')  # Needed for session
app.config['SESSION_TYPE'] = 'filesystem'
app.config['SESSION_FILE_DIR'] = os.path.join(os.path.dirname(__file__), 'flask_session')
app.config['SESSION_PERMANENT'] = False
Session(app)

# Messaggi di sistema e cronologia
messages = [
    {"role": "system", "content": "You are a helpful assistant. Respond to user queries in a friendly manner."}
]

@app.route('/')
def index():
    return '<h2>Chatbot Ollama API</h2><p>Usa POST /chat per interagire con il bot oppure apri <b>chat.html</b> per la chat web.</p>'

# Sostituisci la funzione chat con OpenAI anziché Ollama
@app.route('/chat', methods=['GET', 'POST'])
def chat():
    if request.method == 'GET':
        return '<b>Endpoint API chatbot:</b> invia una richiesta POST con un messaggio JSON {"message": "..."} per ricevere una risposta.'
    try:
        data = request.get_json()
        user_input = data.get('message', '')
        print('DEBUG: user_input:', user_input, file=sys.stderr, flush=True)
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        print('DEBUG: errore parsing json:', tb, file=sys.stderr, flush=True)
        return Response('data: {"response": "Mi dispiace, si è verificato un errore nell\'elaborazione della richiesta. Riprova! Miglioreremo il prodotto."}\n\n', mimetype='text/event-stream')
    import time
    def generate():
        if not user_input:
            yield 'data: {"response": "Non ho ricevuto alcun messaggio. Riprova! Miglioreremo il prodotto."}\n\n'
            return
        start_time = time.time()
        try:
            # Costruisci il contesto in base alla sessione
            chat_messages = [
                {"role": "system", "content": "You are a helpful assistant. Rispondi solo in base al documento caricato, se presente. Se non puoi rispondere, di' che non hai abbastanza informazioni."}
            ]
            doc_content = session.get('doc_content')
            if doc_content:
                chat_messages.append({
                    "role": "system",
                    "content": f"Contenuto documento caricato (usa solo queste informazioni per rispondere):\n{doc_content[:3000]}{'...' if len(doc_content)>3000 else ''}"
                })
            chat_messages.append({"role": "user", "content": user_input})
            print('DEBUG: PRIMA DI OLLAMA', file=sys.stderr, flush=True)
            client = ollama.Client()
            try:
                stream = client.chat(
                    model="llama3",
                    messages=chat_messages,
                    stream=True
                )
            except Exception as e:
                import traceback
                tb = traceback.format_exc()
                yield f'data: {{"response": "Errore nella connessione a Ollama: {str(e)}\n{tb}"}}\n\n'
                return
            bot_response = ""
            chunk_count = 0
            print('DEBUG: INIZIO CICLO STREAM', file=sys.stderr, flush=True)
            for chunk in stream:
                # Timeout di 2 minuti
                if time.time() - start_time > 120:
                    yield 'data: {"response": "Mi dispiace, la risposta ha impiegato troppo tempo. Riprova! Miglioreremo il prodotto."}\n\n'
                    print('DEBUG: TIMEOUT STREAM', file=sys.stderr, flush=True)
                    return
                content = ""
                if isinstance(chunk, dict):
                    if 'message' in chunk and 'content' in chunk['message']:
                        content = chunk['message']['content']
                    elif 'response' in chunk:
                        content = chunk['response']
                elif hasattr(chunk, "message"):
                    msg = chunk.message
                    if isinstance(msg, dict) and 'content' in msg:
                        content = msg['content']
                    elif hasattr(msg, "content"):
                        content = msg.content
                elif hasattr(chunk, "response"):
                    content = chunk.response
                bot_response += content
                chunk_count += 1
                print(f'DEBUG: CHUNK {chunk_count}: {repr(content)}', file=sys.stderr, flush=True)
                yield f"data: {{\"response\": \"{content.replace('\\','\\\\').replace('"','\\"').replace('\n','\\n')}\"}}\n\n"
            print(f'DEBUG: FINE STREAM, chunk_count={chunk_count}', file=sys.stderr, flush=True)
            if chunk_count == 0:
                yield 'data: {"response": "Nessuna risposta ricevuta dal modello AI. Verifica che Ollama sia attivo e il modello disponibile."}\n\n'
                return
            messages.append({"role": "assistant", "content": bot_response})
        except Exception as e:
            import traceback
            tb = traceback.format_exc()
            yield f'data: {{"response": "Mi dispiace, si è verificato un errore interno: {str(e)}\n{tb}"}}\n\n'
    return Response(generate(), mimetype='text/event-stream')

@app.route('/chatweb')
def chatweb():
    # Serve la pagina chat.html dalla directory corrente
    return send_from_directory(os.path.dirname(os.path.abspath(__file__)), 'chat.html')

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'response': 'Nessun file allegato.'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'response': 'Nessun file selezionato.'}), 400
    filename = secure_filename(file.filename)
    ext = os.path.splitext(filename)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        file.save(tmp.name)
        tmp_path = tmp.name
    try:
        doc_content = ''
        if ext in ['.xlsx', '.xls', '.csv']:
            try:
                import pandas as pd
                if ext == '.csv':
                    df = pd.read_csv(tmp_path)
                else:
                    df = pd.read_excel(tmp_path)
                info = f"Excel: {df.shape[0]} righe, {df.shape[1]} colonne. Colonne: {', '.join(df.columns[:5])}{'...' if len(df.columns)>5 else ''}"
                preview = df.head(3).to_markdown(index=False)
                doc_content = df.head(20).to_markdown(index=False)  # Salva più righe per il contesto
                response = f"Ho acquisito il file Excel. {info}\nAnteprima:\n{preview}\n\nDi cosa hai bisogno rispetto a questo documento?"
            except Exception as e:
                response = f"Errore nell'elaborazione del file Excel: {str(e)}"
        elif ext in ['.doc', '.docx']:
            try:
                import docx
                doc = docx.Document(tmp_path)
                text = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
                estratto = text[:700] + ('...' if len(text)>700 else '')
                doc_content = text[:5000]  # Limita la lunghezza
                response = f"Ho acquisito il file Word. Ecco un estratto:\n{estratto}\n\nDi cosa hai bisogno rispetto a questo documento?"
            except Exception as e:
                response = f"Errore nell'elaborazione del file Word: {str(e)}"
        elif ext in ['.pdf']:
            try:
                import PyPDF2
                with open(tmp_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ''
                    for page in reader.pages[:5]:
                        estr = page.extract_text()
                        if estr:
                            text += estr + '\n'
                estratto = text[:700] + ('...' if len(text)>700 else '')
                doc_content = text[:5000]
                response = f"Ho acquisito il PDF. Ecco un estratto:\n{estratto}\n\nDi cosa hai bisogno rispetto a questo documento?"
            except Exception as e:
                response = f"Errore nell'elaborazione del PDF: {str(e)}"
        elif ext in ['.ppt', '.pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                slides = []
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides.append(shape.text.strip())
                estratto = ' | '.join(slides)[:700] + ('...' if len(' | '.join(slides))>700 else '')
                doc_content = ' | '.join(slides)[:5000]
                response = f"Ho acquisito il PowerPoint. Ecco un estratto:\n{estratto}\n\nDi cosa hai bisogno rispetto a questo documento?"
            except Exception as e:
                response = f"Errore nell'elaborazione del PowerPoint: {str(e)}"
        else:
            response = "Formato file non supportato. Sono accettati: Excel, Word, PDF, PowerPoint."
        # Salva il contenuto del documento nella sessione se estratto
        if doc_content:
            session['doc_content'] = doc_content
    except Exception as e:
        response = f"Errore generale nell'elaborazione del file: {str(e)}"
    finally:
        os.remove(tmp_path)
    return jsonify({'response': response})

@app.route('/reset', methods=['POST'])
def reset():
    session.pop('doc_content', None)
    return jsonify({'response': 'Sessione resettata. Puoi caricare un nuovo documento o iniziare una nuova chat.'})

if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5050, debug=True)
